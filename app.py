'''
App made by:
    Eduardo Reyes Alvarez, Ph.D.
Contact:
    eduardo_reyes09@hotmail.com

App version: 
    V06 (Dec 04, 2023): Implemented a way to prepare 1 or 2 pptx in case only one approach (T or FM)
                        was used to quantify, or both. Added progress widgets, and improved dividers.

'''
###################################################################################################

# Import required libraries

import os
import glob
from zipfile import ZipFile
import urllib.request
import pandas as pd
import time
import streamlit as st
from streamlit_option_menu import option_menu

# Python-pptx specific modules

# To make the presentation
from pptx import Presentation
# To specify sizes of images and text (Other options like inches are available)
from pptx.util import Cm, Pt            
# To specify text alignment (the method is used on the text frame, not on the box or the actual text)
from pptx.enum.text import PP_ALIGN     
# To work with line and fill colours
from pptx.enum.dml import MSO_THEME_COLOR

###################################################################################################

# App main layout and page loading

st.set_page_config(
    page_title="Tool 002 - App by Eduardo",
    page_icon=":newspaper:",
    layout="wide")

st.markdown(''' 
            # <span style="color: #bf9000"> PPTX generator for PLA results </span>
            ''', unsafe_allow_html=True)
st.markdown('<hr style="margin-top: +10px; margin-bottom: +10px;">', unsafe_allow_html=True)

# Make a menu of pages on the siderbar, since the app is simple but requires lots of specific details
with st.sidebar:
    selected_page = option_menu("App Menu", ["Generate pptx", "How to use this app", "Info on pptx design"], 
        icons=["filetype-pptx", "patch-question-fill", "columns-gap"], menu_icon="cast", default_index=0)

# Check the selected app page and call the corresponding function to display its content
def change_pages():
    if selected_page == "Generate pptx":
        load_first_page()
    elif selected_page == "How to use this app":
        load_second_page()
    else:
        load_third_page() 

    # Get the working directory for all functions to write files to it, and download the blank pptx template
    st.session_state["current_directory"] = os.path.dirname(os.path.realpath(__file__))
    st.session_state["template_pptx"] = os.path.join(st.session_state["current_directory"], "Template.pptx")
    if not os.path.exists(st.session_state["template_pptx"]):
        template_dir = "https://github.com/EdRey05/Streamlit_projects/raw/main/002_Automated_PPTX_PLA/Template.pptx"
        urllib.request.urlretrieve(template_dir, st.session_state["template_pptx"])

###################################################################################################

# Function to process the input files

@st.cache_data(show_spinner=False)
def process_files(T_and_FM):

    # Extract all the folder structure and files of the zip file provided by the user
    with ZipFile(st.session_state["data_zipfile"], 'r') as zip:
        zip.extractall()

    # The unzipping takes several seconds, so we wait for the folder to appear
    while not "Data" in os.listdir():
        time.sleep(1)

    # Walk through the Data folder to find all csv files
    data_folder = os.path.join(st.session_state["current_directory"], "Data")
    extension_wanted = "*.csv"
    all_csv_files = [file
                    for path, subdir, files in os.walk(data_folder)
                    for file in glob.glob(os.path.join(path, extension_wanted))]

    # Make list with names and paths of the experimental conditions (according to the number of csv files found above)
    exp_conditions_info = []
    
    for csv_file in all_csv_files:
        csv_root_folder = os.path.normpath(csv_file).split(os.path.sep + "Data" + os.path.sep)[1].split(os.path.sep + "Quantification")[0]
        exp_conditions_info.append([csv_root_folder, csv_file])

    # We will store all the information of the images to insert to the ppt here
    all_info_for_slides = []

    # Iterate for each folder contained in the Data folder uploaded
    for condition_info in exp_conditions_info:
        
        # Read the csv file for the current experimental condition
        results_table = pd.read_csv(condition_info[1])

        # Make some edits for easier manipulation and sorting of the subtitles, folders and ROI names
        results_table["Image used"] = results_table["Image used"].str.replace("MAX_","")
        results_table["Image used"] = results_table["Image used"].str.replace(".tif","")
        results_table["Cell quantified"] = results_table["Cell quantified"].str.replace("_1.roi","")
        results_table["Cell quantified"] = results_table["Cell quantified"].apply(int)
        if T_and_FM == "Both" or T_and_FM == "Thresholding only":
            results_table["Particle count threshold"] = results_table["Particle count threshold"].apply(int)
        if T_and_FM == "Both" or T_and_FM == "Find Maxima only":
            results_table["Particle count maxima"] = results_table["Particle count maxima"].apply(int)
        results_table = results_table.sort_values(by=["Image used", "Cell quantified"],ignore_index=True)

        # Iterate through the rows of the csv file (ROIs/cells quantified)
        for _,row in results_table.iterrows():
            
            # Retrieve all the relevant information for that row
            ROI_title = condition_info[0]
            ROI_subtitle = row["Image used"]
            ROI_name = str(row["Cell quantified"])
            ROI_Fluorescence_image = condition_info[1].replace(os.path.join("Quantification", "Results.csv"), 
                                        os.path.join("Cropped cells", "Fluorescence", ROI_subtitle, ROI_name+"_2.jpg"))
            # For now the image and counts will be set to None when only one approach is used (new feature)
            # This allows us to use the same code in the function that makes the slides 
            if T_and_FM == "Thresholding only":
                ROI_Tcount_image = ROI_Fluorescence_image.replace("Fluorescence", "T_Particles").replace("_2.jpg", "_1.jpg")
                ROI_Tcount = row["Particle count threshold"]
                ROI_FMcount_image = None
                ROI_FMcount = None
            elif T_and_FM == "Find Maxima only":
                ROI_Tcount_image = None
                ROI_Tcount = None
                ROI_FMcount_image = ROI_Fluorescence_image.replace("Fluorescence", "FM_Particles").replace("_2.jpg", "_1.jpg")
                ROI_FMcount = row["Particle count maxima"]
            else:
                ROI_Tcount_image = ROI_Fluorescence_image.replace("Fluorescence", "T_Particles").replace("_2.jpg", "_1.jpg")
                ROI_Tcount = row["Particle count threshold"]
                ROI_FMcount_image = ROI_Fluorescence_image.replace("Fluorescence", "FM_Particles").replace("_2.jpg", "_1.jpg")
                ROI_FMcount = row["Particle count maxima"]

            all_info_for_slides.append([ROI_title, ROI_subtitle, ROI_Fluorescence_image, ROI_name, ROI_Tcount_image, ROI_Tcount, ROI_FMcount_image, ROI_FMcount])

    # Prepare empty variables to store the grouped images for each slide
    all_slides_content = []
    temp_slide = []

    # Retrieve the title and subtitle of the very first image to feed the loop at i=0
    current_title = all_info_for_slides[0][0]
    current_subtitle = all_info_for_slides[0][1]

    # Iterate through all the information retrieve for each cell/ROI
    for info in all_info_for_slides:
        
        # Get the current title and subtitle to compare to the reference
        new_title = info[0]
        new_subtitle = info[1]

        # Check the 3 conditions described in text above
        if len(temp_slide)==20 or new_title!=current_title or new_subtitle!=current_subtitle:
            
            # If anything triggers the change of slide, dump the current, empty it and set titles as refs
            all_slides_content.append(temp_slide)
            temp_slide = []
            current_title = new_title
            current_subtitle = new_subtitle
        
        # Always attach the cell/ROI, could be at the same group/slide with others, or to a empty temp slide
        temp_slide.append(info)

    return all_slides_content

###################################################################################################

# Function to generate the presentations and pass the slide maker the content it should insert 

def generate_pptxs(all_slides_content, T_and_FM):

    # This function makes a presentation for the Thresholding and Find Maxima approaches
    # If only one approach was used, this function will make only one presentation

    st.markdown('<hr style="margin-top: +10px; margin-bottom: +10px;">', unsafe_allow_html=True)
    col_1_row_3, col_2_row_3= st.columns([1, 1], gap="medium")

    # Check if we need to generate a presentation for the Thresholding approach
    if T_and_FM == "Both" or T_and_FM == "Thresholding only":

        # Open the presentation uploaded by the user
        presentation_T = Presentation(st.session_state["template_pptx"])

        # Initialize the progress bar
        with col_1_row_3:
            T_progress = st.progress(0, text='Making Thresholding presentation...')

        # Iterate through the image info grouped by slide
        for i,slide_content in enumerate(all_slides_content):

            # Prepare the parameters we need to pass to the function that makes the slides
            current_slide_title = slide_content[0][0]
            current_slide_subtitle = slide_content[0][1]
            image_count_for_slide = len(slide_content)
            F_images_for_slide = [image[2] for image in slide_content]
            F_images_labels = [image[3] for image in slide_content]
            P_images_for_slide = [image[4] for image in slide_content]
            P_images_labels = [image[5] for image in slide_content]

            # Feed the function that makes the slide and inserts the corresponding images
            presentation_T = slide_maker(presentation_T, current_slide_title, current_slide_subtitle, image_count_for_slide, 
                                        F_images_for_slide, P_images_for_slide, F_images_labels, P_images_labels)

            # Update the progress bar
            with col_1_row_3:
                T_progress.progress((i+1)/len(all_slides_content), text=f'Making Thresholding presentation (slide {i+1} of {len(all_slides_content)})')

        # Finally, save this summary presentation after all slides have been created
        presentation_T.save(os.path.join(st.session_state["current_directory"], "Summary_results_T.pptx"))

    #####

    # Check if we need to generate a presentation for the Find Maxima approach
    if T_and_FM == "Both" or T_and_FM == "Find Maxima only":

        # Open the presentation uploaded by the user
        presentation_FM = Presentation(st.session_state["template_pptx"])
        
        # Initialize the progress bar
        with col_2_row_3:
            FM_progress = st.progress(0, text='Making Find Maxima presentation...')

        # Iterate through the image info grouped by slide
        for i,slide_content in enumerate(all_slides_content):

            # Prepare the parameters we need to pass to the function that makes the slides
            current_slide_title = slide_content[0][0]
            current_slide_subtitle = slide_content[0][1]
            image_count_for_slide = len(slide_content)
            F_images_for_slide = [image[2] for image in slide_content]
            F_images_labels = [image[3] for image in slide_content]
            P_images_for_slide = [image[6] for image in slide_content]
            P_images_labels = [image[7] for image in slide_content]
            
            # Feed the function
            presentation_FM = slide_maker(presentation_FM, current_slide_title, current_slide_subtitle, image_count_for_slide, 
                                        F_images_for_slide, P_images_for_slide, F_images_labels, P_images_labels)

            # Update the progress bar
            with col_2_row_3:
                FM_progress.progress((i+1)/len(all_slides_content), text=f'Making Find Maxima presentation (slide {i+1} of {len(all_slides_content)})')

        # Finally, save this summary presentation after all slides have been created
        presentation_FM.save(os.path.join(st.session_state["current_directory"], "Summary_results_FM.pptx"))

###################################################################################################

# Function to add slides to a pptx and inserts images+text 

def slide_maker(presentation_input, current_slide_title, current_slide_subtitle, image_count_for_slide, 
                F_images_for_slide, P_images_for_slide, F_images_labels, P_images_labels):

    # All coordinates are stated always in the same order: From left first, from top second.

    # Title text box dimensions and coordinates (centimeters) 
    title_width = 17
    title_height = 1.5
    title_left_coordinate = 0
    title_top_coordinate = 0
    
    # Subtitle text box dimensions and coordinates (centimeters)
    subtitle_width = 17
    subtitle_height = 1.5
    subtitle_left_coordinate = 17
    subtitle_top_coordinate = 0
    
    # Size and coordinates for the 20 pairs of images (centimeters)
    image_width = 3.25
    image_height = 3
    image_coordinates = [
    (0.25, 2.1, 3.5, 2.1),   (7, 2.1, 10.25, 2.1),   (13.75 , 2.1, 17, 2.1),  (20.5, 2.1, 23.75, 2.1),    (27.25, 2.1, 30.5, 2.1),
    (0.25, 6.4, 3.5, 6.4),   (7, 6.4, 10.25, 6.4),   (13.75, 6.4, 17, 6.4),   (20.5, 6.4, 23.75, 6.4),    (27.25, 6.4, 30.5, 6.4),
    (0.25, 10.7, 3.5, 10.7), (7, 10.7, 10.25, 10.7), (13.75, 10.7, 17, 10.7), (20.5, 10.7, 23.75, 10.7),  (27.25, 10.7, 30.5, 10.7),
    (0.25, 15, 3.5, 15),     (7, 15, 10.25, 15),     (13.75, 15, 17, 15),     (20.5, 15, 23.75, 15),      (27.25, 15, 30.5, 15)
    ]
    
    # Size and coordinates for the 20 pairs of text labels (centimeters) (+3cm top coordinate of images)
    image_labels_width = 3.25
    image_labels_height = 1
    image_labels_coordinates = [
    (0.25, 5.1, 3.5, 5.1),   (7, 5.1, 10.25, 5.1),   (13.75 , 5.1, 17, 5.1),  (20.5, 5.1, 23.75, 5.1),    (27.25, 5.1, 30.5, 5.1),
    (0.25, 9.4, 3.5, 9.4),   (7, 9.4, 10.25, 9.4),   (13.75, 9.4, 17, 9.4),   (20.5, 9.4, 23.75, 9.4),    (27.25, 9.4, 30.5, 9.4),
    (0.25, 13.7, 3.5, 13.7), (7, 13.7, 10.25, 13.7), (13.75, 13.7, 17, 13.7), (20.5, 13.7, 23.75, 13.7),  (27.25, 13.7, 30.5, 13.7),
    (0.25, 18, 3.5, 18),     (7, 18, 10.25, 18),     (13.75, 18, 17, 18),     (20.5, 18, 23.75, 18),      (27.25, 18, 30.5, 18)
    ]
    
    # Create a new slide (layout Blank)
    blank_slide_layout = presentation_input.slide_layouts[6]
    slide = presentation_input.slides.add_slide(blank_slide_layout)
    
    # Make the title for this experimental condition
    left = Cm(title_left_coordinate)
    top = Cm(title_top_coordinate)
    width = Cm(title_width)
    height = Cm(title_height)
    title_textbox = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_textbox.text_frame
    title_text = title_frame.paragraphs[0]
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_text.text = current_slide_title
    title_text.font.bold = True
    title_text.font.size = Pt(32)
    title_text.font.name = "Times New Roman"
    
    # Make the subtitle for the image where the ROI was cropped from
    left = Cm(subtitle_left_coordinate)
    top = Cm(subtitle_top_coordinate)
    width = Cm(subtitle_width)
    height = Cm(subtitle_height)
    subtitle_textbox = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_textbox.text_frame
    subtitle_text = subtitle_frame.paragraphs[0]
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_text.text = current_slide_subtitle
    subtitle_text.font.size = Pt(32)
    subtitle_text.font.name = "Times New Roman"
    
    # Based on the number of images for the current slide, retrieve the neccesary images and coordinates
    for i in range(image_count_for_slide):
        
        # Find the images to insert
        fluorescence_image = F_images_for_slide[i]
        particle_image = P_images_for_slide[i]
        
        # Insert the cropped cell from the Fluorescence folder first
        left = Cm(image_coordinates[i][0])
        top = Cm(image_coordinates[i][1])
        width = Cm(image_width)
        height = Cm(image_height)
        inserting_image = slide.shapes.add_picture(fluorescence_image, left, top, width, height)
        
        # Insert the text label corresponding to the image just inserted above
        left = Cm(image_labels_coordinates[i][0])
        top = Cm(image_labels_coordinates[i][1])
        width = Cm(image_labels_width)
        height = Cm(image_labels_height)
        inserting_image_textbox = slide.shapes.add_textbox(left, top, width, height)
        inserting_image_frame = inserting_image_textbox.text_frame
        inserting_image_text = inserting_image_frame.paragraphs[0]
        inserting_image_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        inserting_image_text.text = F_images_labels[i]
        inserting_image_text.font.size = Pt(20)
        inserting_image_text.font.name = "Times New Roman"
        
        # Insert the cropped cell from the Particles folder second (FM or T particles)
        left = Cm(image_coordinates[i][2])
        top = Cm(image_coordinates[i][3])
        width = Cm(image_width)
        height = Cm(image_height)
        inserting_image2 = slide.shapes.add_picture(particle_image, left, top, width, height)
        inserting_image2.line.fill.solid()
        inserting_image2.line.width = Pt(0.5)
        inserting_image2.line.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    
        # Insert the text label corresponding to the particle counts just inserted above
        left = Cm(image_labels_coordinates[i][2])
        top = Cm(image_labels_coordinates[i][3])
        width = Cm(image_labels_width)
        height = Cm(image_labels_height)
        inserting_image2_textbox = slide.shapes.add_textbox(left, top, width, height)
        inserting_image2_frame = inserting_image2_textbox.text_frame
        inserting_image2_text = inserting_image2_frame.paragraphs[0]
        inserting_image2_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        inserting_image2_text.text = "P="+str(P_images_labels[i])
        inserting_image2_text.font.size = Pt(20)
        inserting_image2_text.font.name = "Times New Roman"
    
    return presentation_input

###################################################################################################

# Function to load the first app page which takes the input Data.zip file and producess the outputs

def load_first_page():
    
    # Create columns for better layout of the buttons 
    col_1_row_1, col_2_row_1 = st.columns([3, 1], gap="large")
    st.markdown('<hr style="margin-top: +15px; margin-bottom: +15px;">', unsafe_allow_html=True)
    col_1_row_2, col_2_row_2, col_3_row_2 = st.columns([1, 1, 1], gap="medium")

    # Show the user a widget to upload the compressed file
    with col_1_row_1:
        uploaded_file = st.file_uploader("Upload compressed file", type=["zip"], accept_multiple_files=False)
    
    # Display a radio button to choose which quantification approach was used
    with col_2_row_1:
        T_and_FM = st.radio(label="Quantification approach used:", 
                            options=["Both","Thresholding only", "Find Maxima only"],
                            index=0,)

    # Display a button so the user decides when to start (in case uploaded the incorrect file)
    with col_1_row_2:
        st.session_state["start_button"] = st.button(label="Generate pptx", type="primary")
    
    # Proceed only when the button to start is pressed and a compressed file has been uploaded
    if st.session_state["start_button"] and uploaded_file:
        
        # Initialize the progress tracker
        with open(uploaded_file.name, "wb") as f:
            f.write(uploaded_file.getvalue())
        st.session_state["data_zipfile"] = os.path.join(os.path.dirname(os.path.abspath(uploaded_file.name)), "Data.zip")

        # Process the files to extract the information needed to import to the slide generator
        all_slides_content = process_files(T_and_FM)
        generate_pptxs(all_slides_content, T_and_FM)
        
        # Check if we need to prepare the path for the Thresholding presentation
        if T_and_FM == "Both" or T_and_FM == "Thresholding only":
            threholding_file_path = os.path.join(st.session_state["current_directory"], "Summary_results_T.pptx")
            st.session_state["path1"] = threholding_file_path

        # Check if we need to prepare the path for the Find Maxima presentation
        if T_and_FM == "Both" or T_and_FM == "Find Maxima only":    
            find_maxima_file_path = os.path.join(st.session_state["current_directory"], "Summary_results_FM.pptx")
            st.session_state["path2"] = find_maxima_file_path
    
    # Show the download button if there is a Thresholding pptx -this way the button persists across reruns
    if "path1" in st.session_state:
        with col_2_row_2:
            st.session_state["download1"] = st.download_button(label="Download Threholding file", 
                               data=open(st.session_state["path1"], "rb").read(), 
                               file_name="Summary_results_T.pptx")

    # Show the download button if there is a Thresholding pptx -this way the button persists across reruns
    if "path2" in st.session_state:
        with col_3_row_2:
            st.session_state["download2"] = st.download_button(label="Download Find Maxima file",  
                                            data=open(st.session_state["path2"], "rb").read(), 
                                            file_name="Summary_results_FM.pptx")

    return

###################################################################################################

# Function to load the second app page with basic instructions on how to use the app

def load_second_page():
    
    st.markdown('''
                ## <span style="color: #0a6640;"> Purpose of this app </span>

                <span style="color: #CE9178">

                This app takes a compressed file with any number of **Proximity Ligation Assay (PLA)** 
                images that were cropped and quantified (see requirements), and prepares a summary 
                Power Point presentation with the results. The slides show for each image:

                <span style="color: #CCCCCC;">
                1. The Fluorescence image used to quantify, with their name below.
                <br><br>
                2. The Particle mask image obtained during the quantification, with their counts below.
                </span>
                <br> <br>

                The idea is that we can quickly preview in one single file all the cell morphologies we 
                picked to analyze, and verify whether the particle counts obtained match what we would 
                expect. This way, we can quickly identify when the selected parameters for quantification 
                produce more/less particles than what we observe in the original fluorescence image. 
                
                By using the PPTX, we can quickly review all the images and quantification results side 
                by side without having to open hundreds of individual images and trackig their names and
                quantification results from the output Excel file.
                
                In addition, we can compare two different quantification approaches (Thresholding vs Find 
                Maxima), and identify possible issues such as:
                
                <span style="color: #CCCCCC;"> - Noisy fluorescent images. </span>

                <span style="color: #CCCCCC;"> - More particles identified than expected.</span>

                <span style="color: #CCCCCC;"> - Less particles identified than expected.</span>

                <span style="color: #CCCCCC;"> - Particles identified seem fine but there are several
                                        black/not coloured (didn't meet criteria to be counted). </span>
                
                These and other problems can be easily identified with the help of the summary PPTX and 
                either lead to the selection of a quantification approach (Thresholding or Find Maxima),
                or prompt the user to repeat the quantiication step with different parameters. The most
                common fixes to these issues include changing: thresholding method, prominence value, 
                particle size + circularity, rolling radius of background subtraction, etc.

                </span>
                ''', unsafe_allow_html=True)
    
    st.divider()
    
    st.markdown('''
                ## <span style="color: #0a6640;"> Requirements </span>

                <span style="color: #CE9178">

                This app was created specifically to generate a result's report for experiments of the
                author's PLA experiments, the only requirement is to upload a compressed zip file. The
                file must not be heavier than 500 MB (if so, you can edit the config.toml file).
                Alternatively, you can make multiple compressed files (all named "Data.zip") and run
                the app for each of them separately.
                
                The quantification is done using the following script for ImageJ/Fiji, which is part 
                of a workflow so it also requires specific image + folder structures:
                [PLA Quantification Script]( 
                https://github.com/EdRey05/Resources_for_Mulligan_Lab/blob/b80eaf75d35665aeb4b7e60ed85685f342d9f125/Tools%20for%20PLA%20quantification/PLA_quantification.py)

                To create the zip file, it is recommended to first create a "Data" folder, then create
                a subfolder for each experimental condition to include, transfer to each subfolder the
                required files and finally compress the "Data" folder into a .zip file. This way, the
                zip file will be called "Data" and contain a "Data" folder, which should contain the
                following:

                </span>
                <span style="color: #CCCCCC;"> 

                1. Any number of subfolders corresponding to each experimental condition to include in
                the PPTX.  

                2. Each condition subfolder has a unique name, but they all should have the same
                content: 1 folder called "Cropped cells", and 1 folder called "Quantification". Both
                of these are produced by the quantification script. 

                3. The "Cropped cells" folder should contain 3 folders: "Fluorescence", "FM_Particles", 
                and "T_Particles". Each of these should have subfolders with the names Row_01_05... to 
                Row_45_51. The "Fluorescence" folder contains ROIs for individual cells with an ID number
                ("Number_2.jpg"). The "FM_Particles" and "T_Particles" folders have the same content of
                a different set of ROIs with an ID number ("Number_1.jpg").

                4. The "Quantification" folder should only have a csv file with the results of the 
                quantification.

                </span>
                <span style="color: #CE9178">

                Here you can find an example zip file that contains the folders and files mentioned above
                for a few experimental conditions of a real PLA experiment (the download button is on the
                right side, between the "Raw" button and the pencil button):
                [Example zip file](
                https://github.com/EdRey05/Resources_for_Mulligan_Lab/blob/caf95fc217cb1c65b4a0b28449c84b35ec10e2fe/Tools%20for%20students/Eduardo%20Reyes/Data.zip)
                
                </span>
                
                ''', unsafe_allow_html=True)
    
    st.divider()

    st.markdown('''
                ## <span style="color: #0a6640;"> Outputs </span>
                
                <span style="color: #CE9178">

                The current version of this notebook produces 2 .pptx files, one using the fluorescence 
                images + the particle mask images prouced by the ***Find Maxima*** approach, and the 
                second one using the fluorescence images + the particle mask images produced by the 
                ***Thresholding*** approach. 
                
                Ideally, we would want to examine both to check which quantification approach works 
                better for the experiment of interest. Under ideal conditions both approaches give very
                similar results. However, the fluorescence image quality is the main factor that 
                influences the results provided by both approaches, which use different principles to
                find the particles.

                For these reasons, we want not only to compare both PPTX summaries, but it is also
                crucial to evaluate consistency in particle detection within each PPTX. Some experimental
                conditions may have more/less background noise or number of particles than other
                conditions, and this may lead to incorrect particle detection and quantification for
                some conditions but not for others. However, the quantification and results report 
                generation are fully automated so that both parts of the workflow can be easily repeated
                to find the most appropriate quantification parameters for each condition and each 
                experiment.

                Furthermore, this app automatically detects experimental conditions and organizes the 
                slides as follows:  
                
                <span style="color: #CCCCCC;">
                -Each experimental condition is shown in a separate slide. Up to 20 pairs of cell images
                can fill each slide. If more slides are needed for a condition, they will all have the 
                same title (top-left).   
                <br><br>
                -If big images containing multiple cells were acquired for each condition, 
                subfolders for each those are created during the quantification. The app identifies
                this and makes multiple slides with the same title, but different subtitle (top-right).
                <br><br>
                <b>NOTE:</b> If you want to see the format of these PPTX presentations, download the 
                example zip file (link above) and run the app with it.
                </span> 

                </span>
                ''', unsafe_allow_html=True)
    return

###################################################################################################

# Function to load the third app page which describes the design and layout of the slides

def load_third_page():
    

    st.markdown('''
            <div style='background-color: #0E6655; padding: 10px; border-radius: 5px; 
                        text-align: center; width: 75%; margin: auto;'>
                <p > Automating the generation of Power Point presentations with our layout is done 
                using the <b>python-pptx</b> library. <br><br>
                For more information on this library, see: 
                <a href="https://pypi.org/project/python-pptx/"> python-pptx (Pypi) </a> </p>
            </div>
                ''', unsafe_allow_html=True)
    st.divider()

    st.markdown('''
                ## <span style="color: #7a3db5;"> Parameters to define </span>

                <span style="color: #CE9178">

                The specific coordinates <b>(in centimeters)</b> for all the desired elements in the
                slides were previously tested by manually arranging images with the intended number
                of rows, columns, images per slide, labels, and text boxes. The size of the elements
                was adjusted to provide sufficient insight into the results while minimizing the 
                number of slides. Once the approximate position of all the content on a single slide
                was set, precise measurements were calculated and applied to ensure reproducibility
                and ease of iteration. 
                
                </span>

                ''', unsafe_allow_html=True)
    
    # Diplay the images of the slide coordinates
    a, img1_container, b = st.columns([0.5, 9, 0.5], gap="small")
    with img1_container:
        st.image(image="https://github.com/EdRey05/Streamlit_projects/raw/main/002_Automated_PPTX_PLA/Automated_PPTX_goal.jpg",
                 caption="Desired slide layout (Title = experimental condition, Subtitle = Image, Image name = ROI name, P = particle count)", 
                 use_column_width=True, )
        
    st.markdown('''
                <span style="color: #CE9178">

                Once the images and text of interest were arranged, the coordinates for each space
                that the objects would fit into were measured as follows. <b>NOTE:</b> That required
                to resize the images and few look stretched in an axis due to their original aspect
                ratio (this can be customized for other experiments depending on the shape and aspect
                ratio of the images).

                </span>

                ''', unsafe_allow_html=True)
    
    c, img2_container, d = st.columns([0.5, 9, 0.5], gap="small")
    with img2_container:
        st.image(image="https://github.com/EdRey05/Streamlit_projects/raw/main/002_Automated_PPTX_PLA/Automated_PPTX_coordinates.jpg",
                 caption="Coordinates for each desired object (x, y)", 
                 use_column_width=True,)
            
    st.markdown('''
                <span style="color: #CE9178">

                All the parameters illustrated above are the following:

                </span>
                <span style="color: #CCCCCC;">

                1. <b>Slide ratio:</b> Use a 16:9 ratio (34cm width, 19cm height - all measurements
                 are given in cm).

                2. <b>Title and subtitle:</b> Two titles at the top, side by side inside a text box
                of 17cm width, and 1.5cm height (right at the top corners). The title has bold Times
                New Roman font, size 32 points, whereas the subtitle has normal Times New Roman font,
                size 32.

                3. <b>Titles separator:</b> Below the title and subtitle text boxes, there is a 0.6cm
                vertical space to the first row of images.

                4. <b>Image size:</b> Each image is resized to 3.25cm width by 3cm height. The pairs
                (yellow + green rectangles) come from the same cell (region of interest). The 
                fluorescence image is on the left side and the particle mask image is on the right
                side, with no space separating them horizontally.

                5. <b>Image labels:</b> There is a 3.25cm width by 1cm height text box right under 
                each image. The text is normal Times New Roman font, size 20 points. The text of the
                fluorescence image indicates the name of the .roi file quantified for that image (so
                we can find it for any purpose). The text of the particle mask image indicates the
                count of particles for that particular region of interest. We should evaluate whether
                the number matches what we see in the fluorescence image, and whether there are
                more/less coloured and black (uncounted) particles in the mask than expected. 

                6. <b>Image pair separation:</b> There are 5 pairs of images per row, 4 pairs per
                column = 20 pairs of original+quantified images (cells or regions of interest). All
                pairs of images are separated by other pairs both horizontally and vertically by 0.25cm.

                7. <b>Image filling order:</b> The images are filled left to right, top-down. They
                appear in the csv file sorted alphabetically, but we inserted them into the slides in
                natural sorting order (more intuitive).

                8. <b>Additional details:</b> Due to the size of the slide and the indicated sizes and
                coordinates, the last row finishes very close from the bottom of the slide, and the 
                last column finishes very close from the right side of the slide (may not be visible as
                the text boxes have no fill and the right side picture has white background).
                
                </span>
                <span style="color: #CE9178">

                **Notes:** 
                
                * Since the fluorescence and particle mask images come from the same cell/region of
                interest, their names are the same but have a different number aft the underscore:
                "100_2.jpg" and "100_1.jpg", respectively.

                * The specific image and label coordinates are hardcoded in tuples (see snippet below).

                * The tuple of coordinates consists of 4 elements: the first is the distance from the
                left edge of the slide for the fluorescence image, the second is the distance from the
                top edge of the slide for the fluorescence image, the third is the distance from the
                left edge of the slide for the particle image, and the fourth is the distance from the
                top edge of the slide for the particle image. This way, we can use the same ROI name to
                get both images by replacing part of the directory and the "_2.jpg" for "_1.jpg",
                however, the index in the for loop will be the same, since we have the coordinates of
                both images in the same element of the list (tuple of 4 coordinates).

                </span>

                ''', unsafe_allow_html=True)
    
    snippet = '''
                # All coordinates are stated always in the same order: From left first, from top second.

                # Title text box dimensions and coordinates (centimeters) 
                title_width = 17
                title_height = 1.5
                title_left_coordinate = 0
                title_top_coordinate = 0
                
                # Subtitle text box dimensions and coordinates (centimeters)
                subtitle_width = 17
                subtitle_height = 1.5
                subtitle_left_coordinate = 17
                subtitle_top_coordinate = 0
                
                # Size and coordinates for the 20 pairs of images (centimeters)
                image_width = 3.25
                image_height = 3
                image_coordinates = [
                (0.25, 2.1, 3.5, 2.1),   (7, 2.1, 10.25, 2.1),   (13.75 , 2.1, 17, 2.1),  (20.5, 2.1, 23.75, 2.1),    (27.25, 2.1, 30.5, 2.1),
                (0.25, 6.4, 3.5, 6.4),   (7, 6.4, 10.25, 6.4),   (13.75, 6.4, 17, 6.4),   (20.5, 6.4, 23.75, 6.4),    (27.25, 6.4, 30.5, 6.4),
                (0.25, 10.7, 3.5, 10.7), (7, 10.7, 10.25, 10.7), (13.75, 10.7, 17, 10.7), (20.5, 10.7, 23.75, 10.7),  (27.25, 10.7, 30.5, 10.7),
                (0.25, 15, 3.5, 15),     (7, 15, 10.25, 15),     (13.75, 15, 17, 15),     (20.5, 15, 23.75, 15),      (27.25, 15, 30.5, 15)
                ]
                
                # Size and coordinates for the 20 pairs of text labels (centimeters) (+3cm top coordinate of images)
                image_labels_width = 3.25
                image_labels_height = 1
                image_labels_coordinates = [
                (0.25, 5.1, 3.5, 5.1),   (7, 5.1, 10.25, 5.1),   (13.75 , 5.1, 17, 5.1),  (20.5, 5.1, 23.75, 5.1),    (27.25, 5.1, 30.5, 5.1),
                (0.25, 9.4, 3.5, 9.4),   (7, 9.4, 10.25, 9.4),   (13.75, 9.4, 17, 9.4),   (20.5, 9.4, 23.75, 9.4),    (27.25, 9.4, 30.5, 9.4),
                (0.25, 13.7, 3.5, 13.7), (7, 13.7, 10.25, 13.7), (13.75, 13.7, 17, 13.7), (20.5, 13.7, 23.75, 13.7),  (27.25, 13.7, 30.5, 13.7),
                (0.25, 18, 3.5, 18),     (7, 18, 10.25, 18),     (13.75, 18, 17, 18),     (20.5, 18, 23.75, 18),      (27.25, 18, 30.5, 18)
                ] 
            '''
    e, snippet_container, f = st.columns([1, 8, 1], gap="small")
    with snippet_container:
        st.code(snippet, language="python")
    
    st.divider()

    st.markdown('''
                ## <span style="color: #7a3db5;"> Overview of the processing strategy </span>

                <span style="color: #CCCCCC;">

                1. Unzip the "Data.zip" file to make a "Data" folder, over which we can iterate/walk 
                through.
                
                2. Once the "Data" folder is ready, get all the folders inside as these are the 
                experimental conditions or groups. Each
                
                3. For each Group folder, there must be a Group/Quantification/results.csv file.

                4. We know that the number of csv files found is the same number of experimental conditions 
                we need to iterate through. We get their paths, names, and we already know the folder 
                structure where the images are located: **Data/ExpCondition/Cropped cells/** , 
                which contains 3 subfolders: **Fluorescence, FM_Particles, and T_Particles**

                5. Once we know the number of experimental condition folders, we will iterate through any 
                number of them to extract the information of all the images.

                6. Since the csv files already contain almost all the information we need (subtitle, ROI 
                names, T particle count and FM particle count), we will iterate through the rows of the 
                csv file instead of walking through the directory of fluorescence images (we could also 
                extract the info from the path but we would need multiple steps to split different 
                sections of the path). This strategy also allows us to easily convert the ROI names into 
                integers so we can sort them properly (natural sorting, the csv file is not in this 
                order).

                </span>
                <span style="color: #CE9178">

                Up to this step, we will have a huge list containing all the info for all the images
                uploaded by the user. The next steps are to group these images to know which ones go 
                together into the same slide, following a few rules:
                
                </span>
                <span style="color: #CCCCCC;">

                * The main idea is to pass 20 images to the function that makes the slides.

                * We start with the first experimental condition (title of slide), first subfolder of 
                the original image (like Row_01_05, which is the subtitle), and then we see how many 
                cells/ROIs were quantified there.

                * If we have exactly 20 (unlikely), we just make one slide. If we have less than 20, we 
                also make one slide and leave empty spots. If we have more than 20 images, we take the 
                first 20, make one slide, take the next (up-to) 20 using the same title and subtitle, 
                make a new slide and so on, until we don't have more images in that subfolder.

                * To accomplish this strategy, we need to iterate through the list generated above 
                (all_info_for_slides), in which all the information for all the images of all the 
                conditions is side by side. Because of that, we can iterate through the same level of 
                the list, take the info of the current item/element, and add its content to a new 
                temporary variable (list) with 20 spots available. Before adding the info of the current 
                item/element, we check whether the 20 spots have been filled, check whether the title 
                has changed, and check whether the subtitle has changed. Any of those 3 cases triggers 
                the immediate jump to a new slide so we have to pass the info of the grouped images
                from the temporary variable to a final variable, clear the temporary variable content
                and then add the current image info to the new slide group. Finally, we repeat this
                process over and over until we have checked all the items/elements, and the final
                variable, product of this loop, will contain all the image information grouped by 
                slide.

                * With the final variable we will be able to iterate through each element (slide), call 
                the function that makes the slides, and pass the current information of the images to
                insert to the slide. We do this over and over until we have made a slide for all groups
                of images specified, and finally we save the presentation. We do this twice, once for
                the images with Thresholding, and once for the images with Find Maxima (info for both
                sets of images and counts is contained in the same variable all_info_for _slides.

                </span>

                ''', unsafe_allow_html=True)
    return

###################################################################################################
change_pages()
###################################################################################################
